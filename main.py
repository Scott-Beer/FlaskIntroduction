# python3 -m streamlit run main.py

import re
import os
import pandas as pd
from pptx import Presentation
import streamlit as st
import datetime

def find_acronyms(text, min_len, max_len):
    """
    Function to find acronyms in a given text.
    Acronyms are defined as alphanumeric characters (case-sensitive, capital letters only), 
    starting and ending with a letter.
    """
    pattern = fr'\b[A-Z][A-Z0-9]{{{min_len-2},{max_len-2}}}[A-Z]\b'
    return re.findall(pattern, text)

def extract_surrounding_words(text, acronym):
    """
    Function to extract 7 words before and after the given acronym.
    """
    words = text.split()
    if acronym in words:
        index = words.index(acronym)
        before = words[max(0, index-7):index]
        after = words[index+1:index+8]
        return ' '.join(before), ' '.join(after)
    return None, None

def process_presentation(file, min_len, max_len):
    """
    Function to process a given PowerPoint presentation,
    returning a DataFrame with acronyms, their contexts and slide numbers.
    """
    presentation = Presentation(file)
    result = []
    for i, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                acronyms = find_acronyms(text, min_len, max_len)
                for acronym in acronyms:
                    before, after = extract_surrounding_words(text, acronym)
                    result.append((acronym, before, after, i+1))
    df = pd.DataFrame(result, columns=['Acronym', 'Before', 'After', 'Slide'])
    return df

def main():
    """
    Main function to handle file upload via streamlit and
    present the resulting DataFrame.
    """
    st.title("PowerPoint Acronym Finder")
    
    st.markdown("""
    ### Instructions:
    1. Use the slider to select the range of acronym lengths you are interested in. Acronyms can be anywhere from 2 to 10 characters long.
    2. Click 'Browse files' to upload a PowerPoint file from your device.
    3. After uploading the file, the application will process it and display a table with all found acronyms, their surrounding context, and the slide number where they were found.
    4. Below the table, you can click 'Download data as CSV' to download the table data for your use.
    5. You'll also see a table showing the count of each acronym found in the file.
    """)
    
    min_len, max_len = st.slider(
        'Select the range of acronym lengths you are interested in',
        min_value=2, max_value=10, value=(2, 7)
    )
    st.markdown("---")  # Add a horizontal line

    uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
    if uploaded_file is not None:
        df = process_presentation(uploaded_file, min_len, max_len)
        
        total_distinct_acronyms = df['Acronym'].nunique()
        st.info(f'Total distinct acronyms: {total_distinct_acronyms}')
        st.markdown("---")  # Add a horizontal line

        st.subheader('Acronym Data:')
        st.dataframe(df)
        
        # Modify the file_name parameter to include the original file name, current date, and 'Acronyms'
        today = datetime.date.today().strftime("%Y%m%d")
        original_filename, _ = os.path.splitext(uploaded_file.name)  # strip extension
        filename = f"{original_filename}_{today}_Acronyms.csv"
        st.download_button(
            "Download data as CSV",
            data=df.to_csv(index=False),
            file_name=filename,
            mime='text/csv'
        )
        st.markdown("---")  # Add a horizontal line

        st.subheader('Acronym Counts:')
        acronym_counts = df['Acronym'].value_counts().reset_index()
        acronym_counts.columns = ['Acronym', 'Count']
        st.dataframe(acronym_counts)

if __name__ == "__main__":
    main()